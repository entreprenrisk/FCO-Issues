from pptx import Presentation
from pptx.util import Inches
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import dataframe_image as dfi
import json
import psycopg2
from psycopg2 import OperationalError
from collections import OrderedDict
import matplotlib as mpl
from PIL import Image
import sys
import copy

import comtypes.client

def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

cmapp = mpl.cm.YlOrRd(np.linspace(0,1,35))
cmapp = mpl.colors.ListedColormap(cmapp[10:,:-1])

cmapn = mpl.cm.PuBuGn(np.linspace(0,1,30))
cmapn = mpl.colors.ListedColormap(cmapn[10:,:-1])

aid_btc = 42840
pids = (1, 2, 3, 4)
pids_crypto = (6, 5, 7, 8)
name_table_res = 'fco_chr_res'
name_columns_res = ['pid', 'time', 'aid', 'signal', 'signal_mean', 'mcluster_pct',
                    'tc_ci5', 'tc_ci95', 't1', 'maxframe', 'tc', '"BubbleSize"']
cols1 = name_columns_res[:-1] + ['bubble size']
cols_name = ['name', 'aid', 'pid', 'ci', 'bd', 'bs', 'bg', 'bp', 'ga', 'tcd', 'sp']
assetTypeCodeToNameMap = {'EQ': 'Single Stocks', 'EQIND': 'Equity Indices', 'BD':'Fixed Income',
                          'CMD': 'Commodities', 'EX': 'Forex', 'CRP': 'CryptoX'}
NUM_RESULTS = 3

PPT_DEFAULT_SIZE = [10.833, 7.5]
LAYOUT_IMG, LAYOUT_TEXT, LAYOUT_LIST= 2, 9, 10
IMG_CHARTTYP_LOC = [Inches(0), Inches(0), Inches(10.8), Inches(7.2)]
#ToDo type {0: IMAGE, 1: TEXT, 2:TABLE}
prs=Presentation("FCO.pptx")

def get_slide_center_to_image(path):
    im = Image.open(path)
    width = im.width / im.info['dpi'][0]
    height = im.height / im.info['dpi'][1]
    left = Inches(PPT_DEFAULT_SIZE[0] - width) / 2
    top = Inches(PPT_DEFAULT_SIZE[1] - height) / 2
    return [left, top]

def add_ppt_slide(type, obj):
    if type == 0:
        lyt=prs.slide_layouts[LAYOUT_IMG] # choosing a slide layout
        slide=prs.slides.add_slide(lyt) # adding a slide
        slide.shapes.add_picture(obj['name'], *obj['loc'])
        if 'title' in obj:
            title = slide.shapes.title
            title.text= obj['title']
    if type == 1 and any(key in obj for key in ['title', 'subtitle']):
        print(obj)
        lyt = prs.slide_layouts[LAYOUT_TEXT]
        slide = prs.slides.add_slide(lyt)
        if 'title' in obj:
            title = slide.shapes.title
            title.text= obj['title']
        if 'subtitle' in obj:
            subtitle = slide.placeholders[1]
            subtitle.text=obj['subtitle']
    if type == 2 and all(key in obj for key in ['title', 'list']):
        lyt = prs.slide_layouts[LAYOUT_LIST]
        slide = prs.slides.add_slide(lyt)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        p = tf
        if 'title' in obj:
            title_shape.text = obj['title']
        if 'list' in obj:
            for ix,item in enumerate(obj['list']):
                p.text = str(item)
                if ix < len(obj['list'])-1:
                    p = tf.add_paragraph()
    if type == 3 and all(key in obj for key in ['path', 'index']):
        # copy from external presentation all objects into the existing presentation
        external_pres = Presentation(obj['path'])

        # specify the slide you want to copy the contents from
        ext_slide = external_pres.slides[obj['index']]

        slide_layout = None
        if 'layout' in obj:
            slide_layout = prs.slide_layouts[obj['layout']]
        else:
            for inx, layouts in enumerate(prs.slide_layouts):
                if layouts.name == ext_slide.slide_layout.name:
                    slide_layout = layouts
                    break
        if slide_layout:
            # create now slide, to copy contents to
            curr_slide = prs.slides.add_slide(slide_layout)

            # now copy contents from external slide, but do not copy slide properties
            # e.g. slide layouts, etc., because these would produce errors, as diplicate
            # entries might be generated

            for shp in ext_slide.shapes:
                el = shp.element
                newel = copy.deepcopy(el)
                curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

def move_ppt_slide(src, to=None):
    last = prs.slides.length -1
    prs.move_slide(src, last if to is None else to)

def save_ppt(name):
    prs.save("%s.pptx"%name) # saving file

def create_connection(db_name, db_user, db_password, db_host, db_port):
    connection = None
    try:
        connection = psycopg2.connect(
            database=db_name,
            user=db_user,
            password=db_password,
            host=db_host,
            port=db_port,
        )
        print("Connection to PostgreSQL DB successful")
    except OperationalError as e:
        print(f"The error '{e}' occurred")
    return connection

def execute_read_query(connection, query):
    cursor = connection.cursor()
    result = None
    try:
        cursor.execute(query)
        result = cursor.fetchall()
    except OperationalError as e:
        print(f"The error '{e}' occurred")
    return result

def get_table(connection,name_table,name_columns='*',where=None,n_limit=None):
    if name_columns!='*':
        name_columns = ','.join(name_columns)
    query_get_table = "SELECT %s FROM %s"%(name_columns,name_table)
    if where is not None:
        query_get_table += " WHERE (%s)"%where
    if n_limit is not None:
        query_get_table += ' LIMIT %d'%n_limit
    if name_columns=='*':
        query_get_table_columns = "SELECT column_name FROM INFORMATION_SCHEMA.COLUMNS WHERE TABLE_NAME = '%s'"%name_table
        cols = execute_read_query(connection,query_get_table_columns)
        name_columns = [col[0] for col in cols]
    else:
        name_columns = name_columns.split(',')
    data = execute_read_query(connection,query_get_table)
    df = pd.DataFrame(data,columns=name_columns)
    return df

def get_table_sql(connection, sql):
    cursor = connection.cursor()
    df = pd.DataFrame()
    try:
        cursor.execute(sql)
        df = pd.DataFrame(cursor.fetchall())
        df.columns = [desc[0] for desc in cursor.description]
    except OperationalError as e:
        print(f"The error '{e}' occurred")
    return df

def parse(x):
    try:
        x = x.astype(int)
    except:
        pass
    return x

def parse_round(x):
    try:
        x = np.round(x,2)
    except:
        pass
    return x

def get_first_tc(x):
    return x[0]

def parse_tc_date(x):
    return str(x.date())

def parse_pid_to_name(x):
    if x==1 or x==6:
        x = 'Short Term'
    elif x==2 or x==5:
        x = 'Super Short'
    elif x==3 or x==7:
        x = 'Medium Term'
    elif x==4 or x==8:
        x = 'Long Term'
    else:
        pass
    return x

def query_price_indicator(connection, aid, ts, te):
    query_indicator_pos = "select time,signal_mean,pid from fco_chr_res where (aid=%d and time between '%s' and '%s' and \"BubbleSize\">0) order by 1" % (
    aid, ts, te)
    query_indicator_neg = "select time,signal_mean,pid from fco_chr_res where (aid=%d and time between '%s' and '%s' and \"BubbleSize\"<0) order by 1" % (
    aid, ts, te)
    query_price = "select time,adjclose from fco_prices  where (aid=%d and time between '%s' and '%s') order by 1" % (
    aid, ts, te)

    dfi_pos = pd.DataFrame(execute_read_query(connection, query_indicator_pos), columns=['t', 'ci', 'pid']).set_index(
        't')
    dfi_pos = dfi_pos.pivot(columns='pid').T.reset_index(drop=True).reindex(range(4)).T
    cols = [dfi_pos.columns[i] for i in [1, 0, 2, 3]]
    dfi_pos = dfi_pos[cols].fillna(0)
    dfi_pos.columns = ['PB CI SS', 'PB CI S', 'PB CI M', 'PB CI L']

    dfi_neg = pd.DataFrame(execute_read_query(connection, query_indicator_neg), columns=['t', 'ci', 'pid']).set_index(
        't')
    dfi_neg = dfi_neg.pivot(columns='pid').T.reset_index(drop=True).reindex(range(4)).T
    cols = [dfi_neg.columns[i] for i in [1, 0, 2, 3]]
    dfi_neg = dfi_neg[cols].fillna(0)
    dfi_neg.columns = ['NB CI SS', 'NB CI S', 'NB CI M', 'NB CI L']

    dfp = pd.DataFrame(execute_read_query(connection, query_price), columns=['t', 'p']).set_index('t')
    dfi_pos = 100 * dfi_pos.reindex(dfp.index).fillna(0)
    dfi_neg = -100 * dfi_neg.reindex(dfp.index).fillna(0)
    dfi_pos.index.name = None
    dfi_neg.index.name = None
    dfp.index.name = None

    return dfi_pos.astype(np.float64), dfi_neg.astype(np.float64), dfp.astype(np.float64)

def plot_price_in(dict):
    aid, name, type, geo, ts, t2, lppls = dict['aid'], dict['name'], dict['type'], dict['geog'], dict['ts'], dict['t2'], dict['lppls']
    conn = connection if type != 'CRP' else connection_crypto
    dfi_pos, dfi_neg, dfp = query_price_indicator(conn, aid, ts, t2)

    if name in ['Venezuela Se General', 'Jpm Venezuela Reer Cpi (2010=100)']:
        return

    save_name = type +  '_' + geo + '_' + str(aid)
    print(save_name)

    # plot

    fig, ax = plt.subplots(1, 1, figsize=(12, 8))


    color = ['b', 'g', 'y', 'r']
    ax1 = ax.twinx()
    ax.set_facecolor('#e5e5e5')

    dfp.plot(ax=ax, c='k', legend=False, grid=True)

    if not dfi_pos.replace(0, np.nan).dropna(how='all').empty:
        dfi_pos.plot(ax=ax1, linewidth=.5, grid=False, color=color, stacked=False)
        dfi_pos.plot.area(ax=ax1, alpha=0.4, linewidth=0., grid=False, color=color, stacked=False)
    if not dfi_neg.replace(0, np.nan).dropna(how='all').empty:
        dfi_neg.plot(ax=ax1, linewidth=.5, grid=False, color=color, stacked=False)
        dfi_neg.plot.area(ax=ax1, alpha=0.4, linewidth=0., grid=False, color=color, stacked=False)
    h, l = ax1.get_legend_handles_labels()
    ax1.legend(h[:4], ['CI SS', 'CI S', 'CI M', 'CI L'], ncol=2, frameon=False, loc=1)

    ax.set_xlim(dfp.index[0], dfp.index[-1])
    ax1.set_xlim(dfp.index[0], dfp.index[-1])
    ax1.set_ylim(-100, 100)
    ax.set_ylabel('Price')
    ax1.set_ylabel('Indicator [%]')
    if lppls:
        tc= pd.Timestamp(dict['tcd'])
        bd = dict['bd']
        if aid != aid_btc:

            t1 = dfp.loc[:t2].iloc[-int(bd):].index[0]
            leg = False
            if pd.Timestamp(ts) + 5 * pd.Timedelta('1d') < tc < pd.Timestamp(t2):
                ax.axvline(tc, linewidth=1., color='k', linestyle='--',
                           label=r'Latest Mean Critical Time $\bar{t}_c$ Estimate')
                leg = leg or True
            if pd.Timestamp(ts) + 5 * pd.Timedelta('1d') < t1 < pd.Timestamp(t2):
                ax.axvline(t1, linewidth=1., color='g', linestyle='--',
                           label=r'Latest Bubble Start Time $t_1$ Estimate')
                leg = leg or True
            if leg:
                h, l = ax.get_legend_handles_labels()
                ax.legend(h[1:], l[1:], ncol=1, frameon=False, loc=2)

        ispos = dict['ispos']
        color = 'g' if ispos else 'r'
        save_name = save_name+ '_' + ('0_pb' if ispos else '1_nb')
        plt.suptitle('%s[%s  %s] %s' %('[%s] '%(geo if geo else ''), 'POS' if ispos else 'NEG', dict['rank'], name) , fontsize=12, fontweight="bold", color=color) #PB/NB DS LPPLS Confidence Multi-Scale Indicators for

        plt.title(', '.join([r'DS LPPLS $\bf\ {}$ CI $\bf\ {}$%'.format(dict['pid'], dict['ci']),
                             r'Size $\bf\ {}$%'.format(dict['bs']),
                             r'CAGR $\bf\ {}$%'.format(dict['bg']),
                             r'Duration $\bf\ {}$ days'.format(bd), '\n'
                             r'Progress $\bf\ {}$%'.format(dict['bp']),
                             r'Geometric Mean $\bf\ {}$'.format(dict['ga']),
                             r'Critical Time $\bf\ {}$ [$t_c$]'.format(tc.strftime('%Y.%m.%d')),
                             r'Scenario Prob $\bf\ {}$%'.format(dict['sp'])]), fontsize=10)
    else:
        plt.title('%s' % name, fontsize=12)

    plt.tight_layout()
    pic_name = 'indicator_plots/' + save_name + '.png'
    fig.savefig(pic_name, dpi=400)
    plt.close(fig)
    add_ppt_slide(0, {'name':pic_name, 'loc':IMG_CHARTTYP_LOC})

sql = ''' select a.type, r.time, count(distinct a.aid) as total, COUNT(distinct a.aid) FILTER (WHERE signal and "BubbleSize" >= 0.1 and signal_mean >= 0.1) as pos, 
            -1 * COUNT(distinct a.aid) FILTER (WHERE signal and "BubbleSize" <= -0.1 and signal_mean >= 0.1) as neg
          from fco_chr_res r
          join asset_base_data2 a on r.aid=a.aid
          and time between '{}' and '{}'
          group by 1,2
          order by 1,2 '''

def get_bubble_fraction_plot(t2s, t2e):
    df_frac = get_table_sql(connection, sql.format(t2s, t2e))
    df_frac['time'] = df_frac['time'].apply(pd.Timestamp)
    df_frac.set_index(['time', 'type'], inplace=True)
    df_frac['fracp'] = 100 * df_frac['pos'] / df_frac['total']
    df_frac['fracn'] = 100 * df_frac['neg'] / df_frac['total']
    df_frac = df_frac.round({'fracp': 0, 'fracn': 0})

    types = list(assetTypeCodeToNameMap.keys())
    if not include_crypto:
        types.remove('CRP')
    iterables = [pd.date_range(pd.Timestamp(t2s), pd.Timestamp(t2e), freq='B'), types]
    ddf = pd.DataFrame(index=pd.MultiIndex.from_product(iterables, names=["time", "type"]))
    df_frac = pd.concat([ddf, df_frac], axis=1)

    fig, ax = plt.subplots(1, 1, figsize=(12, 7))
    ax.set_facecolor('#e5e5e5')

    # df_frac.unstack(level=1).plot(subplots=True)
    colors = ['C%d' % i for i in range(10)]
    colors[3] = 'g'
    colors[2] = 'C3'
    colors[4] = 'C2'

    kgrp = df_frac.reset_index('type').groupby(['type'])
    for col in ['fracp', 'fracn']:
        ix = 1
        for key, grp in kgrp:
            grp[col].plot(ax=ax, fontsize=12, alpha=0.9, color=colors[ix], label=assetTypeCodeToNameMap[key], legend=False)
            ix = ix + 1

    h, l = ax.get_legend_handles_labels()
    ax.legend(h[len(types):], l[len(types):], frameon=False, ncol=6, loc=3)

    title = 'Fraction of Positive / Negative Bubble Signals for different Asset Classes'
    ax.set_title(title, fontsize=14)
    ax.set_ylabel('Bubble Fraction [%]', fontsize=12)

    ax.xaxis.grid(linestyle='--', which='major', color='white')
    ax.yaxis.grid(linestyle='--', which='major', color='white')

    plt.tight_layout()
    pic_path = 'indicator_plots/bubble_overview_202103.png'
    fig.savefig(pic_path, dpi=500)

    add_ppt_slide(0, {'name': pic_path, 'loc': IMG_CHARTTYP_LOC,
                      'title': 'General Results - The Big Picture'})

def main(t2s, t2e):
    global connection
    global connection_crypto
    global include_crypto

    df_geog = get_table(connection, 'fco_asset_geog').set_index('num')
    df_geog.loc[df_geog[df_geog.name == 'United States'].index, 'continent'] = 'United States'
    df_geog = df_geog.replace('Americas', 'America (excl. US)').sort_index()

    df_pars = get_table(connection, 'fco_chr_params').set_index('pid').T

    # standard assets from fcodb
    name_columns = ['id', 'symbol', 'name', 'code', 'geog', 'isactive']
    df_assets = get_table(connection, 'fco_assets', name_columns, n_limit=None).set_index('id')
    df_assets = df_assets.reindex(df_assets['isactive'].replace(False, np.nan).dropna().index)

    df_assets['name'] = [str(n).title() for n in df_assets['name']]
    df_assets['geog'] = df_geog.reindex(df_assets.geog)['continent'].values

    # cryptoassets from fcocrypto db
    if include_crypto:
        df_assets_crypto = get_table(connection_crypto, 'fco_assets', name_columns, n_limit=None).set_index('id')
        df_assets_crypto = df_assets_crypto.reindex(df_assets_crypto['isactive'].replace(False, np.nan).dropna().index)
        df_assets_crypto['geog'] = 'World'

        # merge
        df_assets = pd.concat([df_assets, df_assets_crypto], axis=0)

    df_assets.rename(columns={'code': 'type'}, inplace=True)
    assert len(df_assets.index.drop_duplicates()) == len(df_assets)

    fraction_aids = []
    get_bubble_fraction_plot(t2s, t2e)
    fraction_aids = get_bubble_fraction_table(df_assets, t2s, t2e)

    add_ppt_slide(1, {'title': 'Appendix'})
    df_assets_ = df_assets.reset_index().rename(columns={'id':'aid'})
    for aid in [4710 , 4769 , 6661, 4713, 43]:
        if aid not in fraction_aids:
            dict = df_assets_[df_assets_['aid'] == aid].to_dict('records')[0]
            dict.update({'lppls': False, 'ts':t2s, 't2':t2e})
            plot_price_in(dict)

    save_ppt('FCO_Test')

def get_bubble_fraction_table(df_assets, t2s, t2e):
    for t2 in pd.date_range(t2s, t2e, freq='B')[-1:]:
        print(t2.date(), t2e, end='\r')

        t2 = str(t2.date())

        where = "pid in %s and time = '%s' " % (pids, t2)  # and signal and signal_mean>0.1 and \"BubbleSize\"<-0.1"%t2
        df_res = get_table(connection, name_table_res, name_columns_res, where)

        if include_crypto:
            where_crypto = "pid in %s and time = '%s' " % (pids_crypto, t2)
            df_res_crypto = get_table(connection_crypto, name_table_res, name_columns_res, where_crypto)
            df_res_crypto['pid'] = df_res_crypto['pid'].replace(dict(zip(pids_crypto, pids)))

            df_res = pd.concat([df_res, df_res_crypto], axis=0)
        df_res.columns = cols1

        df_res = pd.concat([df_res.reset_index(), df_assets.reindex(df_res.aid).reset_index(drop=True)], axis=1)
        df_res = df_res[df_res.isactive == True]
        del df_res['index']

        df_res_unique_aid = df_res.drop_duplicates(['aid'])
        n_tot = pd.DataFrame()
        for ix in df_res_unique_aid['type'].value_counts().index:
            n_i = df_res_unique_aid[df_res_unique_aid.type == ix]['geog'].value_counts()
            n_tot = pd.concat([n_tot, n_i.to_frame(ix)], axis=1, sort=True).sort_index(ascending=True, axis=1)


        res_neg = df_res[df_res['bubble size'] < 0]
        res_neg['signal'] = (res_neg['signal'] == True) & (res_neg['signal_mean'] >= 0.1) & (res_neg['bubble size'] <= -0.1)
        res_pos = df_res[df_res['bubble size'] > 0]
        res_pos['signal'] = (res_pos['signal'] == True) & (res_pos['signal_mean'] >= 0.1) & (res_pos['bubble size'] >= 0.1)

        # count all assets with valid positive / negative bubble signals
        # compute bubble fractions

        n_neg = res_neg[res_neg.signal].drop_duplicates('aid')
        n_tot_neg = pd.DataFrame()
        for ix in n_neg['type'].value_counts().index:
            n_i = n_neg[n_neg.type == ix]['geog'].value_counts()
            n_tot_neg = pd.concat([n_tot_neg, n_i.to_frame(ix)], axis=1, sort=True).sort_index(ascending=True, axis=1)
        n_tot_neg = n_tot_neg.fillna(0) * pd.DataFrame(1, n_tot.index, n_tot.columns)

        n_pos = res_pos[res_pos.signal].drop_duplicates('aid')
        n_tot_pos = pd.DataFrame()
        for ix in n_pos['type'].value_counts().index:
            n_i = n_pos[n_pos.type == ix]['geog'].value_counts()
            n_tot_pos = pd.concat([n_tot_pos, n_i.to_frame(ix)], axis=1, sort=True).sort_index(ascending=True, axis=1)
        n_tot_pos = n_tot_pos.fillna(0) * pd.DataFrame(1, n_tot.index, n_tot.columns)

        # select asset class - geog - pairs to consider
        tuples = [('BD', ''),
         ('CMD', ''),
         ('EQ', ''), ('EQ', 'Europe'), ('EQ', 'United States'),
         ('EQIND', ''), ('EQIND', 'Europe'), ('EQIND', 'United States'), ('EQIND', 'World'),
         ('EX', '')] + ([('CRP', '')] if include_crypto else [])
        mix = pd.MultiIndex.from_tuples(tuples, names=['Asset Class', 'Region'])

        col_names = ['Analyzed Assets', 'Fraction of Pos. Bubbles [%]', 'Fraction of Neg. Bubbles [%]']
        res = pd.DataFrame(np.nan, index=mix, columns=col_names)

        # asset classes for which to compute the total fractions ix = (Type, Geog)
        for ix in res.index:
            if not ix[-1]:
                # assign entire asset class values
                res.loc[ix, res.columns[0]] = n_tot.sum().loc[ix[0]]
                res.loc[ix, res.columns[1]] = 100 * n_tot_pos.sum().loc[ix[0]] / n_tot.sum().loc[ix[0]]
                res.loc[ix, res.columns[2]] = 100 * n_tot_neg.sum().loc[ix[0]] / n_tot.sum().loc[ix[0]]
            else:
                ixi = (ix[-1], ix[0])
                res.loc[ix, res.columns[0]] = n_tot.loc[ixi]
                res.loc[ix, res.columns[1]] = 100 * n_tot_pos.loc[ixi] / n_tot.loc[ixi]
                res.loc[ix, res.columns[2]] = 100 * n_tot_neg.loc[ixi] / n_tot.loc[ixi]

        # postprocessing / styling
        res = res.fillna(0).round(0)  # .replace('World','Global')
        ixa = [rix[0] for rix in res.index]
        for i in range(1, len(ixa)):
            if ixa[i] in ixa[:i]:
                ixa[i] = ''
        res['Region'] = [rix[1] for rix in res.index]
        res = res[[res.columns[i] for i in [-1, 0, 1, 2]]]
        res.index = [assetTypeCodeToNameMap['BD'], assetTypeCodeToNameMap['CMD'], assetTypeCodeToNameMap['EQ'], '', '',
                     assetTypeCodeToNameMap['EQIND'], '', '', '', assetTypeCodeToNameMap['EX']]+ (assetTypeCodeToNameMap['CRP'] if include_crypto else [])
        res.iloc[-3, 0] = 'Global'
        res = res.apply(parse)
        # if not RES:
        bbl_fr_tb = 'indicator_plots/bubble_fraction_tbl.png'
        dfi.export(res, bbl_fr_tb)
        add_ppt_slide(0, {'name':bbl_fr_tb, 'loc':get_slide_center_to_image(bbl_fr_tb), 'title': 'General Results as of %s'%t2})
        return print_signal_table(df_assets.type.drop_duplicates().values.tolist(),
                               df_assets.geog.drop_duplicates().values.tolist(),
                                n_pos, n_neg, tuples, t2)

def compute_rel_quant(ni):
    # get tc and tc as date
    ni['tc'] = ni['tc'].apply(get_first_tc)
    ni['tc'] = ni['tc'] * ni['maxframe']
    ni['tcd'] = (ni['time'] + (-ni['maxframe'] + ni['tc'].round(0)) * pd.Timedelta('1d'))
    ni['tcd'] = ni['tcd'].apply(parse_tc_date)

    # get bubble cagr, duration, progress and geometric average
    ni['bd'] = ni['maxframe'] - ni['t1']
    # ni['bp'] = (ni['tc'] - ni['t1']) / ni['bd'] # maxframe is t2 (or t2+1, not sure)
    ni['bp'] = ni['bd'] / (ni['tc'] - ni['t1'])  # maxframe is t2 (or t2+1, not sure)
    ni['bg'] = np.exp((np.log(ni['bubble size'] + 1) * 252 / ni['bd'])) - 1
    ni['ga'] = (np.abs(ni['signal_mean'] * ni['bg'] * ni['bp']) ** (1. / 3))

    # filter out the ones with 0.25 <= bp and tc-t2>=20 business days
    # (i.e. bubbles with tc too far in the future or too long ago)
    ni = ni[(0.25 <= ni['bp']) & (ni['bp'] <= 1.0)]  # ((ni['tc']-ni['maxframe'])>=-20)]

    # ['name', 'aid', 'pid', 'signal_mean', 'bd', 'bubble size', 'bg', 'bp', 'ga', 'tcd', 'mcluster_pct']
    # ['name', 'aid', 'pid', 'ci', 'bd', 'bs', 'bg', 'bp', 'ga', 'tcd', 'sp']
    for col in ['signal_mean', 'bubble size', 'bg', 'bp', 'ga', 'mcluster_pct']:
        ni[col] *= 100
    ni['pid'] = ni['pid'].apply(parse_pid_to_name)

    # select for each asset class - region - pair the three largest signals in terms of geometric average
    ni_ga_max = ni.reindex(ni.groupby('aid').ga.idxmax())
    ni_ga_max['tg'] = ni_ga_max['type'] + ' ' + ni_ga_max['geog']
    ix = [ixi[1] if hasattr(ixi, "__getitem__") else ixi for ixi in ni_ga_max.groupby('tg').ga.nlargest(3).index]
    ni_ga_max = ni_ga_max.loc[ix]
    return ni_ga_max

def comp_top_signal_rankedby_ga(n_pos, n_neg):
    cols_sel = ['name', 'aid', 'pid', 'signal_mean', 'bd', 'bubble size', 'bg', 'bp', 'ga', 'tcd', 'mcluster_pct']

    # cols_pid = ['ci1','ci2','ci3','ci4']
    # for both bubble types, compute relevant quantities
    ress = OrderedDict()
    for i, ni in enumerate([n_pos, n_neg]):
        if not ni.empty:
            ni_ga_max = compute_rel_quant(ni)
            sigi = OrderedDict()
            for ty in ni_ga_max.type.drop_duplicates():
                sigi[ty] = OrderedDict()
                nt = ni_ga_max[ni_ga_max.type == ty]
                for ge in nt.geog.drop_duplicates():

                    ntg = nt[nt.geog == ge][cols_sel]
                    ntg.columns = cols_name

                    ntg = ntg.drop_duplicates('name', keep='first').reset_index(drop=True)
                    ntg.index += 1
                    ntg = ntg.apply(parse_round)

                    sigi[ty][ge] = ntg

            # compute overall top signals ranked by ga for each class
            for ty in sigi:

                sty = sigi[ty]

                s_all = pd.DataFrame()
                for ge in sty:
                    if ge == 'Overall':
                        continue
                    s_all = pd.concat([s_all, sty[ge]], axis=0)

                s_all = s_all.sort_values('ga', ascending=False).iloc[:NUM_RESULTS]

                sigi[ty]['Overall'] = s_all

            ress[i] = sigi
    return ress

def print_signal_table(n_pos, n_neg, tyge, t2):

    result_aids = set()
    ### PREPARE AND PRINT REPORT SIGNAL TABLES

    ress = comp_top_signal_rankedby_ga(n_pos, n_neg)

    ts = str((pd.Timestamp(t2) - 365 * pd.Timedelta('1d')).date())[:-2] + '01'
    conc = []
    php = [pd.DataFrame([len(cols_name) * ['']], index=[r'%s Bubbles'%('Negative' if i else 'Positive')], columns=cols_name) for i in [0,1]]
    for ix, (ty, ge) in enumerate(tyge):
        print(ty, ge)
        for iix,i in enumerate(ress.keys()):
            first = True if iix == 0 and (ix == 0 or ty != tyge[ix - 1][0]) else False
            if ty in ress[i] and ge in ress[i][ty]:
                if first:
                    add_ppt_slide(1, {'title':'%s %s'%(assetTypeCodeToNameMap[ty], ge)})
                stygep = ress[i][ty][ge].reset_index(drop=True)
                for ni in ['Venezuela Se General', 'Jpm Venezuela Reer Cpi (2010=100)']:
                    if ni in stygep.name.tolist():
                        stygep = stygep.drop(stygep.loc[stygep.name == ni].index[0])
                        stygep = stygep.reset_index(drop=True)
                stygep.index += 1
                conc.extend([php[i], stygep])

                styge = ress[i][ty][ge]

                aids = styge.aid.tolist()
                result_aids.update(aids)
                if ty == 'CRP' and ge == 'Overall':
                    aids += [aid_btc]
                for j, aid in enumerate(aids):
                    dict = styge[styge['aid'] == aid].to_dict('records')[0]
                    dict.update({'type':ty, 'geog':ge, 'ispos': i==0, 'lppls':True, 'ts':ts, 't2':t2, 'rank':j})
                    plot_price_in(dict)
            else:
                if first:
                    add_ppt_slide(1, {'title': assetTypeCodeToNameMap[ty], **({'subtitle': 'No bubbles to report'} if ix == len(tyge) -1 or ty != tyge[ix+1][0] else {}) })
    return result_aids

if __name__ == "__main__":
    connection = psycopg2.connect('<DB_URL_1>')
    include_crypto = False
    connection_crypto = psycopg2.connect('<DB_URL_2>') if include_crypto else None
    main('2020-02-01', '2021-02-28')
